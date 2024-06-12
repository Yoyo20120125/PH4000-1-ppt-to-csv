from time import sleep
from pyautogui import press, hotkey
from pyperclip import copy, paste
from pptx import Presentation
import os

FILE_NAME = "U27"

file_path = str(os.path.dirname(os.path.abspath(__file__))) + "\\" + FILE_NAME + ".pptx"
print("The file path is: " + file_path)


def to_ppt():
    ppt = Presentation(file_path)
    with open(file_path[:-5] + ".txt", 'w', encoding = 'utf-8') as f:
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    f.write(shape.text)

def read():
    num = ["0", "1", "2", "3", "4", "5", "6", "7", "8", "9"]
    text = ""
    word = []
    chinese = []

    with open("C:\\Users\\QRV4k\\Desktop\\" + FILE_NAME + ".txt", "r", encoding = "utf-8") as file:
        text = file.readlines()

    for i in range(len(text)):
        if text[i][0] in num:
            t = text[i].replace("\t", " ").split(" ")
            line = []
            for i in t:
                if i != "":
                    line.append(i)
            word.append(line[1])
            chinese_text = line[2] + line[3]
            print(chinese_text)
            chinese.append(chinese_text.strip())

def type():
    print("5 second!")
    sleep(5)
    for i in range(len(word)):
        copy(word[i])
        hotkey('ctrl', 'v')
        press("Enter")
    press("Right")
    for i in range(20):
        press("Up")
    for i in range(len(chinese)):
        copy(chinese[i])
        hotkey('ctrl', 'v')
        press("Enter")

to_ppt()
read()
type()